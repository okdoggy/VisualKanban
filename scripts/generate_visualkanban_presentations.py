from pathlib import Path
from typing import Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path("docs/presentations")
ASSET_DIR = OUT_DIR / "assets"
OUT_DIR.mkdir(parents=True, exist_ok=True)

FONT = "Malgun Gothic"
COLOR_BG = RGBColor(247, 249, 253)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK = RGBColor(28, 33, 43)
COLOR_MUTED = RGBColor(92, 103, 122)
COLOR_PRIMARY = RGBColor(30, 88, 199)
COLOR_SECONDARY = RGBColor(32, 160, 118)
COLOR_ACCENT = RGBColor(121, 82, 179)


def set_bg(slide, color=COLOR_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_text(run, size, bold=False, color=COLOR_DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_top_band(slide, label: str):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.48))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()
    tf = band.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Visual Kanban  |  {label}"
    p.alignment = PP_ALIGN.LEFT
    style_text(p.runs[0], 13, True, COLOR_WHITE)


def add_title(slide, title: str, subtitle: Optional[str] = None):
    tbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.62), Inches(12.1), Inches(0.8))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    style_text(p.runs[0], 30, True)
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.67), Inches(1.18), Inches(12.0), Inches(0.6))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        style_text(sp.runs[0], 14, False, COLOR_MUTED)


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: Iterable[str],
    border_color=COLOR_PRIMARY,
):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = border_color
    box.line.width = Pt(1.8)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    p0 = tf.paragraphs[0]
    p0.text = title
    style_text(p0.runs[0], 18, True)
    for line in lines:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        style_text(p.runs[0], 13, False, COLOR_DARK)
        p.space_before = Pt(4)


def add_image_card(slide, image: Path, x: float, y: float, w: float, h: float, caption: Optional[str] = None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = RGBColor(210, 219, 236)
    card.line.width = Pt(1.5)

    if image.exists():
        slide.shapes.add_picture(str(image), Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(h - 0.5))
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(h - 0.5))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(236, 240, 248)
        ph.line.color.rgb = RGBColor(205, 214, 230)
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"이미지 없음: {image.name}"
        p.alignment = PP_ALIGN.CENTER
        style_text(p.runs[0], 13, True, COLOR_MUTED)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if caption:
        cap = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + h - 0.32), Inches(w - 0.24), Inches(0.2))
        ctf = cap.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.text = caption
        style_text(cp.runs[0], 10, False, COLOR_MUTED)


def add_kpi(slide, x: float, y: float, w: float, h: float, label: str, value: str, tone=COLOR_PRIMARY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = tone
    card.line.width = Pt(1.8)
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = label
    style_text(p1.runs[0], 11, False, COLOR_MUTED)
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    style_text(p2.runs[0], 24, True, tone)
    p2.alignment = PP_ALIGN.CENTER


def add_flow_boxes(slide, y: float, labels: List[str]):
    count = len(labels)
    total_w = 11.9
    box_w = 2.1
    gap = (total_w - (count * box_w)) / (count - 1 if count > 1 else 1)
    x = 0.7
    boxes = []
    for i, label in enumerate(labels):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_WHITE
        box.line.color.rgb = COLOR_PRIMARY if i % 2 == 0 else COLOR_SECONDARY
        box.line.width = Pt(1.6)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        style_text(p.runs[0], 13, True)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        boxes.append(box)
        x += box_w + gap

    for i in range(len(boxes) - 1):
        a, b = boxes[i], boxes[i + 1]
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            a.left + a.width,
            a.top + a.height // 2,
            b.left,
            b.top + b.height // 2,
        )
        line.line.color.rgb = COLOR_MUTED
        line.line.width = Pt(1.3)


def add_chart(slide, x, y, w, h, categories, series: Dict[str, List[int]]):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_major_gridlines = False


def add_table(slide, x, y, w, h, rows, cols, headers, data):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    table.first_row = True
    for c in range(cols):
        table.columns[c].width = Inches(w / cols)
    for c, title in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = title
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 238, 252)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style_text(p.runs[0], 12, True, COLOR_DARK)
    for r_idx, row in enumerate(data, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx != 1 else PP_ALIGN.LEFT
            style_text(p.runs[0], 11, False, COLOR_DARK)


def build_intro_deck(out_file: Path):
    prs = Presentation()

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    add_top_band(slide, "소개 기능 중심 발표")
    add_title(slide, "Visual Kanban", "협업 운영을 하나의 화면으로 통합한 실무형 플랫폼")
    add_image_card(slide, ASSET_DIR / "dashboard.png", 0.7, 1.65, 7.9, 5.2, "실제 서비스 화면: Dashboard")
    add_bullets(
        slide,
        8.75,
        1.65,
        3.85,
        5.2,
        "이번 발표 포인트",
        [
            "페이지별 기능을 실제 화면과 함께 설명",
            "사용 시나리오 중심으로 가치 전달",
            "도입 시 운영효과를 수치/표로 정리",
        ],
        border_color=COLOR_ACCENT,
    )

    # 2 agenda + KPI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "소개 기능 중심 발표")
    add_title(slide, "발표 구성", "기능 설명 + 운영 방식 + 기대효과")
    add_bullets(
        slide,
        0.7,
        1.6,
        7.1,
        3.6,
        "Agenda",
        [
            "1) 화면 구조와 페이지별 주요 기능",
            "2) 권한/프로젝트 운영 방식",
            "3) 실제 업무 흐름 예시",
            "4) 도입 기대효과",
        ],
    )
    add_kpi(slide, 8.1, 1.9, 1.35, 1.4, "핵심 페이지", "7개", COLOR_PRIMARY)
    add_kpi(slide, 9.62, 1.9, 1.35, 1.4, "권한 등급", "3+1", COLOR_SECONDARY)
    add_kpi(slide, 11.14, 1.9, 1.35, 1.4, "배포 방식", "Docker", COLOR_ACCENT)
    add_image_card(slide, ASSET_DIR / "search.png", 8.1, 3.55, 4.4, 2.8, "Search: To do/칸반/간트 통합 검색")

    # 3 information architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "소개 기능 중심 발표")
    add_title(slide, "전체 화면 구조", "워크스페이스 + 현재화면 + 페이지 콘텐츠")
    add_flow_boxes(slide, 2.0, ["좌측 Workspace", "상단 Current View", "중앙 Page"])
    add_bullets(
        slide,
        0.7,
        3.45,
        12.0,
        2.9,
        "구조 포인트",
        [
            "워크스페이스 접기/펼치기를 지원해 집중 모드 전환 가능",
            "검색, 계정, 스타일/언어 설정이 상단/좌측에 일관 배치",
            "Dashboard, To do, Whiteboard, Kanban, Gantt, 사용자 관리로 이동",
        ],
    )

    # 4 dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "페이지별 기능")
    add_title(slide, "Dashboard", "오늘 해야 할 일과 프로젝트 현황을 한 번에 파악")
    add_image_card(slide, ASSET_DIR / "dashboard.png", 0.7, 1.6, 8.1, 5.7, "프로젝트/할 일/진행상태를 카드형으로 요약")
    add_bullets(
        slide,
        8.95,
        1.6,
        3.7,
        5.7,
        "핵심 기능",
        [
            "최근 프로젝트 자동 유지 (계정별)",
            "칸반/간트 빠른 진입 버튼",
            "팀 진행상태 스냅샷 확인",
            "스탠드업/주간회의 시작 화면으로 활용",
        ],
    )

    # 5 todo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "페이지별 기능")
    add_title(slide, "To do", "개인 실행력 관리 영역")
    add_image_card(slide, ASSET_DIR / "todo.png", 0.7, 1.6, 8.1, 5.7, "개인 일정·반복·우선순위 관리")
    add_bullets(
        slide,
        8.95,
        1.6,
        3.7,
        5.7,
        "실무 활용",
        [
            "개인 루틴과 팀 태스크를 분리 운영",
            "반복/완료 기준으로 누락 방지",
            "우선순위 기준으로 당일 집중 목록 정리",
        ],
        border_color=COLOR_SECONDARY,
    )

    # 6 kanban
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "페이지별 기능")
    add_title(slide, "Kanban Board", "상태 기반 협업 보드")
    add_image_card(slide, ASSET_DIR / "kanban.png", 0.7, 1.6, 8.1, 5.7, "Backlog / To do / In Progress / Done")
    add_bullets(
        slide,
        8.95,
        1.6,
        3.7,
        5.7,
        "핵심 포인트",
        [
            "카드 이동 중심의 빠른 상태 업데이트",
            "상세 팝업 심플화 (제목/설명/담당/우선순위)",
            "하이라이트로 중요 항목 가독성 강화",
            "담당자 자동완성 입력으로 대규모 사용자 대응",
        ],
        border_color=COLOR_ACCENT,
    )

    # 7 gantt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "페이지별 기능")
    add_title(slide, "Gantt Chart", "일정·계층·진행을 동시에 관리")
    add_image_card(slide, ASSET_DIR / "gantt.png", 0.7, 1.6, 8.1, 5.7, "트리 + 타임라인 동시 편집")
    add_bullets(
        slide,
        8.95,
        1.6,
        3.7,
        5.7,
        "핵심 포인트",
        [
            "행 드래그로 순서/계층 구조 변경",
            "그래프 드래그로 일정 이동 및 조정",
            "상세 팝업 단순화로 핵심 속성 집중 편집",
            "하이라이트 시 트리/그래프 모두 동일 강조",
        ],
        border_color=COLOR_PRIMARY,
    )

    # 8 whiteboard + user mgmt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "페이지별 기능")
    add_title(slide, "Whiteboard / 사용자 관리", "아이디어 공유 + 프로젝트 인원 운영")
    add_image_card(slide, ASSET_DIR / "whiteboard.png", 0.7, 1.6, 6.1, 5.7, "화이트보드 협업 화면")
    add_image_card(slide, ASSET_DIR / "admin-users.png", 6.95, 1.6, 5.7, 5.7, "사용자/프로젝트 멤버 관리")

    # 9 role table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "운영 정책")
    add_title(slide, "권한 운영 모델", "역할 + 프로젝트 참여 여부로 기능 접근 제어")
    add_table(
        slide,
        0.7,
        1.65,
        12.0,
        3.0,
        rows=5,
        cols=4,
        headers=["구분", "설명", "프로젝트 관리", "칸반/간트/화이트보드 작성"],
        data=[
            ["Owner", "프로젝트 생성자/최고 권한", "가능", "가능"],
            ["Write", "협업 편집 권한", "가능(멤버 추가/삭제 포함)", "가능"],
            ["Read", "열람 중심 권한", "불가", "불가(조회만)"],
            ["Admin", "서비스 전체 관리자", "가능", "가능"],
        ],
    )
    add_bullets(
        slide,
        0.7,
        4.9,
        12.0,
        1.7,
        "운영 팁",
        [
            "프로젝트 참여자만 삭제 가능하도록 제한해 안전성 확보",
            "대규모 사용자를 고려해 멤버 추가는 자동완성 입력 방식 채택",
        ],
        border_color=COLOR_SECONDARY,
    )

    # 10 scenario flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "운영 시나리오")
    add_title(slide, "실제 업무 진행 시나리오", "요청 접수부터 완료 보고까지")
    add_flow_boxes(slide, 2.0, ["요청 등록", "Kanban 진행", "Gantt 일정조정", "Whiteboard 협의", "완료/공유"])
    add_bullets(
        slide,
        0.7,
        3.45,
        12.0,
        2.7,
        "핵심 메시지",
        [
            "업무 상태, 일정, 논의 흔적이 같은 프로젝트 컨텍스트에 유지",
            "리더/실무자/참여자가 동일 데이터로 커뮤니케이션",
            "툴 전환 비용을 줄이고 완료 리드타임을 단축",
        ],
    )

    # 11 benefits chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "도입 기대효과")
    add_title(slide, "팀 운영 지표 개선 목표", "※ 아래 수치는 내부 운영 목표(예시)입니다")
    add_chart(
        slide,
        0.8,
        1.65,
        7.3,
        4.9,
        categories=["가시성", "처리속도", "일정정합성", "협업만족도"],
        series={
            "도입 전": [48, 52, 46, 55],
            "도입 후 목표": [82, 78, 80, 84],
        },
    )
    add_bullets(
        slide,
        8.35,
        1.8,
        4.25,
        4.7,
        "예상 효과",
        [
            "현황 파악 시간 단축",
            "우선순위 기반 처리 속도 향상",
            "일정 충돌 조기 인지",
            "팀 간 전달 누락 감소",
        ],
        border_color=COLOR_ACCENT,
    )

    # 12 closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    add_top_band(slide, "마무리")
    add_title(slide, "정리", "Visual Kanban은 실무 협업을 '하나의 운영 화면'으로 단순화합니다")
    add_bullets(
        slide,
        0.8,
        1.8,
        12.0,
        4.4,
        "결론",
        [
            "기능 측면: 할 일/칸반/간트/화이트보드/사용자 관리를 일관 UX로 통합",
            "운영 측면: 권한/프로젝트 모델로 안전하고 명확한 협업 구조 제공",
            "기술 측면: Docker + PostgreSQL 기반으로 실제 팀 운영에 바로 투입 가능",
            "다음 단계: 팀별 운영 규칙 템플릿 정의 및 온보딩 진행",
        ],
    )
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(12.0), Inches(0.6))
    tf = thanks.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Q&A"
    p.alignment = PP_ALIGN.CENTER
    style_text(p.runs[0], 24, True, COLOR_PRIMARY)

    prs.save(out_file)


def build_engineering_deck(out_file: Path):
    prs = Presentation()

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "Visual Kanban Engineering Deck", "기술 스택 · 아키텍처 · 워크플로우")
    add_image_card(slide, ASSET_DIR / "kanban.png", 0.7, 1.65, 6.1, 5.2, "Frontend 실사용 화면")
    add_image_card(slide, ASSET_DIR / "gantt.png", 6.95, 1.65, 5.7, 5.2, "복잡 인터랙션 화면")

    # 2 architecture overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "아키텍처 개요", "Next.js 단일 앱 + API + PostgreSQL")
    add_flow_boxes(slide, 1.95, ["Browser", "Next.js App Router", "API /api/state", "PostgreSQL"])
    add_bullets(
        slide,
        0.7,
        3.35,
        8.1,
        2.9,
        "설계 의도",
        [
            "프론트/백엔드 통합 배포로 운영 복잡도 감소",
            "클라이언트 상태 + 서버 저장 동기화로 다중 사용자 대응",
            "버전 충돌(Optimistic Concurrency)로 동시 편집 안전성 확보",
        ],
    )
    add_image_card(slide, ASSET_DIR / "dashboard.png", 8.95, 3.35, 3.75, 2.9, "실제 운영 화면")

    # 3 tech stack matrix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "기술 스택 매트릭스")
    add_table(
        slide,
        0.7,
        1.65,
        12.0,
        4.5,
        rows=8,
        cols=3,
        headers=["레이어", "기술", "적용 이유"],
        data=[
            ["UI Framework", "Next.js 16 / React 19", "App Router 기반 구조화 + 생산성"],
            ["Language", "TypeScript", "도메인 모델/권한 로직의 타입 안정성"],
            ["State", "Zustand + persist", "가벼운 전역 상태 + 클라이언트 지속성"],
            ["DnD", "dnd-kit", "칸반/간트 드래그 UX 구현"],
            ["Whiteboard", "Excalidraw", "실시간 협업 메모/스케치 경험"],
            ["Database", "PostgreSQL + pg", "기업 환경 적합, 신뢰성 높은 트랜잭션/JSONB"],
            ["Infra", "Docker Compose", "앱+DB 일괄 배포/복구 단순화"],
        ],
    )

    # 4 data & permission model
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "데이터/권한 모델")
    add_flow_boxes(slide, 1.95, ["User", "Project", "Membership", "Task/Todo/Scene"])
    add_bullets(
        slide,
        0.7,
        3.35,
        8.1,
        2.8,
        "권한 해석 규칙",
        [
            "BaseRole(admin/editor/viewer) + ProjectMembership(owner/write/read) 조합",
            "feature 단위 resolveRole → canRead/canWrite로 최종 권한 판정",
            "프로젝트 참여자가 아니면 작성 권한 자동 제한",
        ],
        border_color=COLOR_SECONDARY,
    )
    add_image_card(slide, ASSET_DIR / "admin-users.png", 8.95, 3.35, 3.75, 2.8, "사용자/멤버십 관리 화면")

    # 5 auth flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "워크플로우 A: 인증/세션")
    add_flow_boxes(slide, 2.0, ["로그인 입력", "계정 검증", "쿠키 발급(vk_user)", "초기비번 변경", "앱 라우팅"])
    add_bullets(
        slide,
        0.7,
        3.45,
        12.0,
        2.7,
        "구현 포인트",
        [
            "없는 계정 + 0000 입력 시 파트 입력 후 신규 계정 생성",
            "초기 비밀번호(0000) 강제 변경 후 실제 화면 접근 허용",
            "AppShell에서 세션 상태를 기준으로 라우팅 가드 처리",
        ],
    )

    # 6 sync flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "워크플로우 B: 상태 동기화")
    add_flow_boxes(slide, 1.95, ["hydrate", "GET /api/state", "로컬 편집", "PUT expectedVersion", "409 충돌 재동기화"])
    add_bullets(
        slide,
        0.7,
        3.35,
        8.1,
        2.8,
        "안정화 포인트",
        [
            "seedRevision 기반 데이터 마이그레이션/초기화 제어",
            "세션 사용자 누락 스냅샷 방어 로직으로 로그인 튕김 방지",
            "pending snapshot 상태에서는 폴링 적용을 지연해 데이터 역전 방지",
        ],
        border_color=COLOR_ACCENT,
    )
    add_image_card(slide, ASSET_DIR / "search.png", 8.95, 3.35, 3.75, 2.8, "검색 페이지(동기화 반영 확인)")

    # 7 interaction implementation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "워크플로우 C: 복잡 인터랙션 구현")
    add_image_card(slide, ASSET_DIR / "kanban.png", 0.7, 1.65, 6.0, 4.9, "Kanban - 카드 이동 / 상세 편집")
    add_image_card(slide, ASSET_DIR / "gantt.png", 6.95, 1.65, 5.7, 4.9, "Gantt - 트리/그래프 드래그")
    add_bullets(
        slide,
        0.7,
        6.0,
        12.0,
        1.2,
        "공통 패턴: 사용자 입력 → 권한체크 → Store 업데이트 → 동기화 API 반영",
        [],
    )

    # 8 reliability chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "품질 게이트 & 운영 안정성", "※ 지표는 운영 체크리스트 기준 예시")
    add_chart(
        slide,
        0.8,
        1.65,
        7.3,
        4.9,
        categories=["Typecheck", "Lint", "Build", "권한회귀", "동기화회귀"],
        series={
            "기준 통과율 목표": [100, 100, 100, 98, 98],
            "현재 점검 수준": [100, 100, 100, 96, 95],
        },
    )
    add_bullets(
        slide,
        8.35,
        1.8,
        4.25,
        4.7,
        "검증 방식",
        [
            "npm run typecheck",
            "npm run lint",
            "npm run build",
            "로그인/권한/동기화 회귀 시나리오 점검",
        ],
    )

    # 9 docker deployment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "배포 토폴로지 (Docker)")
    add_flow_boxes(slide, 2.0, ["User Browser", "visualkanban container", "postgres container", "volume(backups)"])
    add_bullets(
        slide,
        0.7,
        3.45,
        12.0,
        2.8,
        "운영 포인트",
        [
            "DATABASE_URL, sync 옵션으로 환경별 동작 제어",
            "DB 컨테이너 healthcheck 완료 후 앱 컨테이너 기동",
            "볼륨 기반 데이터 영속화 + 정기 백업 정책 권장",
        ],
        border_color=COLOR_SECONDARY,
    )

    # 10 roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "개발자 설명")
    add_title(slide, "확장 로드맵")
    add_table(
        slide,
        0.7,
        1.65,
        12.0,
        3.5,
        rows=5,
        cols=3,
        headers=["우선순위", "항목", "목표"],
        data=[
            ["P1", "실시간 알림/멘션", "변경 인지 속도 향상"],
            ["P1", "감사로그 확장", "변경 추적/보안 대응 강화"],
            ["P2", "템플릿/자동화 규칙", "반복 업무 최소화"],
            ["P2", "운영 대시보드", "장애/성능 가시성 강화"],
        ],
    )
    add_bullets(
        slide,
        0.7,
        5.35,
        12.0,
        1.35,
        "핵심 원칙: 기능 확장보다도 데이터 정합성·권한 안전성·운영 단순성을 우선",
        [],
        border_color=COLOR_ACCENT,
    )

    # 11 closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    add_top_band(slide, "마무리")
    add_title(slide, "개발 관점 요약")
    add_bullets(
        slide,
        0.8,
        1.8,
        12.0,
        4.6,
        "정리",
        [
            "Visual Kanban은 실무 협업 UX와 운영 안전성을 동시에 고려한 구조입니다.",
            "PostgreSQL 동기화 + 권한 모델 + Docker 배포로 즉시 서비스 가능한 기반을 갖췄습니다.",
            "앞으로는 관측성/알림/자동화 레이어를 추가해 운영 성숙도를 높일 계획입니다.",
        ],
    )
    qna = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(12.0), Inches(0.6))
    tf = qna.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Technical Q&A"
    p.alignment = PP_ALIGN.CENTER
    style_text(p.runs[0], 22, True, COLOR_PRIMARY)

    prs.save(out_file)


def main():
    intro = OUT_DIR / "VisualKanban_소개_페이지상세.pptx"
    eng = OUT_DIR / "VisualKanban_기술_워크플로우_개발자용.pptx"
    build_intro_deck(intro)
    build_engineering_deck(eng)
    print(intro)
    print(eng)


if __name__ == "__main__":
    main()
